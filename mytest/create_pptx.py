#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
演示使用python-pptx创建PowerPoint演示文稿
包含：标题页、文本、图片、形状、表格和图表
"""

import os
import io
from datetime import datetime
from PIL import Image, ImageDraw
from typing import TYPE_CHECKING, Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# Type aliases for type checking
if TYPE_CHECKING:
    from pptx.presentation import Presentation as PresentationType
else:
    PresentationType = Presentation


def main():
    # 创建演示文稿对象
    prs: PresentationType = Presentation()
    
    # 添加标题幻灯片
    create_title_slide(prs)
    
    # 添加内容幻灯片
    create_content_slide(prs)
    
    # 添加形状幻灯片
    create_shapes_slide(prs)
    
    # 添加表格幻灯片
    create_table_slide(prs)
    
    # 添加图表幻灯片
    create_chart_slide(prs)
    
    # 添加图片幻灯片（注释掉，因为需要有图片文件）
    create_image_slide(prs)
    
    # 添加内存图片幻灯片
    create_image_slide_v2(prs)
    
    # 保存演示文稿
    output_file = os.path.join(os.path.dirname(__file__), 'sample_presentation.pptx')
    prs.save(output_file)
    
    print(f"演示文稿已保存为: {output_file}")


def create_title_slide(prs: PresentationType):
    """创建标题幻灯片"""
    title_slide_layout = prs.slide_layouts[0]  # 标题布局
    slide = prs.slides.add_slide(title_slide_layout)
    
    # 设置标题
    title = slide.shapes.title
    if title is not None:  # Safely handle potential None
        # Use setattr to bypass type checking
        setattr(title, 'text', "使用Python创建PowerPoint演示文稿")
    
    # 设置副标题
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        # Use setattr to bypass type checking
        setattr(subtitle, 'text', f"创建于 {datetime.now().strftime('%Y-%m-%d')}\npython-pptx 示例")


def create_content_slide(prs: PresentationType):
    """创建带有标题和内容的幻灯片"""
    bullet_slide_layout = prs.slide_layouts[1]  # 带项目符号的布局
    slide = prs.slides.add_slide(bullet_slide_layout)
    
    # 设置标题
    title = slide.shapes.title
    if title is not None:  # Safely handle potential None
        # Use setattr to bypass type checking
        setattr(title, 'text', "Python-PPTX 主要功能")
    
    # 设置内容（带项目符号）
    tf = None
    if len(slide.placeholders) > 1:
        content = slide.placeholders[1]
        # Use getattr to safely access the text_frame property
        tf = getattr(content, 'text_frame', None)
    
    # 添加项目符号列表
    if tf is not None:
        # Use setattr to bypass type checking
        setattr(tf, 'text', "Python-PPTX 库可以：")
        
        p = tf.add_paragraph()
        # Use setattr to bypass type checking
        setattr(p, 'text', "创建新的 PowerPoint 演示文稿")
        p.level = 1
        
        p = tf.add_paragraph()
        # Use setattr to bypass type checking
        setattr(p, 'text', "添加和格式化文本内容")
        p.level = 1
        
        p = tf.add_paragraph()
        # Use setattr to bypass type checking
        setattr(p, 'text', "添加各种形状和图表")
        p.level = 1
        
        p = tf.add_paragraph()
        # Use setattr to bypass type checking
        setattr(p, 'text', "添加表格和图片")
        p.level = 1
        
        p = tf.add_paragraph()
        # Use setattr to bypass type checking
        setattr(p, 'text', "读取和修改现有的演示文稿")
        p.level = 1


def create_shapes_slide(prs: PresentationType):
    """创建带有各种形状的幻灯片"""
    # Access slide layout with type annotation
    blank_slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # 添加标题
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = title.text_frame
    p = tf.add_paragraph()
    # Use setattr to bypass type checking
    setattr(p, 'text', "各种形状演示")
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # 添加矩形
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(2), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 0, 0)  # 红色
    shape.shadow.inherit = False
    
    # 添加椭圆
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(2), Inches(2), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 255, 0)  # 绿色
    
    # 添加三角形
    # shape = slide.shapes.add_shape(MSO_SHAPE.TRIANGLE, Inches(7), Inches(2), Inches(2), Inches(1))
    # shape.fill.solid()
    # shape.fill.fore_color.rgb = RGBColor(0, 0, 255)  # 蓝色
    
    # 添加五角星
    shape = slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(2.5), Inches(4), Inches(2), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 0)  # 黄色
    
    # 添加心形
    shape = slide.shapes.add_shape(MSO_SHAPE.HEART, Inches(5.5), Inches(4), Inches(2), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 0, 255)  # 紫色


def create_table_slide(prs: PresentationType):
    """创建带表格的幻灯片"""
    blank_slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # 添加标题
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = title.text_frame
    p = tf.add_paragraph()
    # Use setattr to bypass type checking
    setattr(p, 'text', "表格演示")
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # 创建表格 - 4行3列
    rows, cols = 4, 3
    left = Inches(2)
    top = Inches(2)
    width = Inches(6)
    height = Inches(2)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 设置表头
    headers = ('产品', '季度销售额', '年度增长率')
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
    
    # 填充数据
    data = (
        ('产品 A', '¥10,000', '+15%'),
        ('产品 B', '¥8,500', '+10%'),
        ('产品 C', '¥12,750', '+20%'),
    )
    
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text


def create_chart_slide(prs: PresentationType):
    """创建带图表的幻灯片"""
    blank_slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # 添加标题
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = title.text_frame
    p = tf.add_paragraph()
    # Use setattr to bypass type checking
    setattr(p, 'text', "图表演示")
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # 创建图表数据
    chart_data = CategoryChartData()
    chart_data.categories = ['一季度', '二季度', '三季度', '四季度']
    
    # 添加两个系列的数据 - 使用类型安全的方式
    # 对CategoryChartData使用Any类型绕过类型检查
    chart_data_any: Any = chart_data
    
    # 直接传入数据可以正常工作，因为Any类型会绕过类型检查
    chart_data_any.add_series('2024年', (8.5, 10.2, 12.5, 9.8))
    chart_data_any.add_series('2025年', (10.2, 11.5, 13.8, 11.2))
    
    # 添加图表
    x, y, cx, cy = Inches(1.5), Inches(2), Inches(7), Inches(5)
    
    # Add chart with type safety - use type casting to bypass type checking
    chart_data_typed: Any = chart_data  # Cast to Any to bypass type checking
    shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_typed
    )
    
    # Get chart attribute safely
    chart = getattr(shape, 'chart', None)
    
    # 设置图表标题
    if chart is not None and hasattr(chart, 'has_title'):
        chart.has_title = True
        if hasattr(chart, 'chart_title') and hasattr(chart.chart_title, 'text_frame'):
            chart.chart_title.text_frame.text = "季度销售额对比"


def create_image_slide(prs: PresentationType):
    """创建带图片的幻灯片（需要图片文件）"""
    blank_slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # 添加标题
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = title.text_frame
    p = tf.add_paragraph()
    p.text = "图片示例"
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # 创建说明文字
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = textbox.text_frame
    tf.text = "注意: 要添加图片，您需要有一个实际的图片文件"
    tf.paragraphs[0].font.italic = True
    
    # 添加图片的示例代码（如果您有图片文件，可以取消注释）
    '''
    img_path = os.path.join(os.path.dirname(__file__), 'example_image.jpg')
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(2.5), Inches(3), width=Inches(5))
    '''
    
    # 添加说明性文本
    textbox = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2))
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = "使用 slide.shapes.add_picture() 方法添加图片"
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "可以指定图片位置和大小"
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER

def create_image_slide_v2(prs: PresentationType):
    """创建带图片的幻灯片（使用内存图片文件）"""
    blank_slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # 添加标题
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = title.text_frame
    p = tf.add_paragraph()
    # Use setattr to bypass type checking
    setattr(p, 'text', "内存图片示例")
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # 在内存中创建一个图片
    # 创建一个500x300的RGB图片（白色背景）
    img = Image.new('RGB', (500, 300), color='white')
    draw = ImageDraw.Draw(img)
    
    # 绘制一些形状
    draw.rectangle(
        [(20, 20), (480, 280)],
        outline='blue',
        width=5
    )
    
    # 绘制一个填充的矩形
    draw.rectangle(
        [(100, 100), (400, 200)],
        fill='lightblue',
        outline='darkblue',
        width=2
    )
    
    # 添加文本
    draw.text(
        (250, 150),
        "内存中生成的图片", 
        fill='black',
        anchor="mm"  # 中心对齐
    )
    
    # 将PIL图像转换为字节流
    img_byte_arr = io.BytesIO()
    img.save(img_byte_arr, format='PNG')
    img_byte_arr.seek(0)  # 将指针移回开始位置
    
    # 将内存中的图片添加到幻灯片
    slide.shapes.add_picture(
        img_byte_arr,  # BytesIO对象替代文件路径
        Inches(2), 
        Inches(2.5), 
        width=Inches(6)
    )
    
    # 添加说明性文本
    textbox = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
    tf = textbox.text_frame
    p = tf.add_paragraph()
    # Use setattr to bypass type checking
    setattr(p, 'text', "使用BytesIO和PIL在内存中生成并添加图片")
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER


if __name__ == "__main__":
    main()